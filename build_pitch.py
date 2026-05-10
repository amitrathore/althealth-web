from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ──────────────────────────────────────────────
INK        = RGBColor(0x14, 0x1c, 0x26)
MUTED      = RGBColor(0x4a, 0x56, 0x68)
BRAND      = RGBColor(0x1a, 0x3d, 0xb8)
BRAND_DEEP = RGBColor(0x0c, 0x1a, 0x4f)
ACCENT     = RGBColor(0xb8, 0x89, 0x2e)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG   = RGBColor(0xf7, 0xf3, 0xeb)
CARD_BG    = RGBColor(0xff, 0xff, 0xff)
SOFT_BLUE  = RGBColor(0xe8, 0xef, 0xff)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        from pptx.util import Pt as pt
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Manrope", italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def eyebrow(slide, text, x, y, w):
    add_text(slide, text.upper(), x, y, w, Inches(0.35),
             size=9, bold=True, color=BRAND, font="Manrope")

def heading(slide, text, x, y, w, size=36, color=INK):
    add_text(slide, text, x, y, w, Inches(1.4),
             size=size, bold=True, color=color, font="Fraunces")

def body(slide, text, x, y, w, h=None, size=14, color=MUTED):
    add_text(slide, text, x, y, w, h or Inches(0.5),
             size=size, color=color, font="Manrope")

def bullet_box(slide, items, x, y, w, h, size=13, color=MUTED):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {item}"
        run.font.name  = "Manrope"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        p.space_before = Pt(4)

def dark_slide_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=BRAND_DEEP)

def light_slide_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=LIGHT_BG)

def accent_bar(slide):
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=ACCENT)

def slide_number(slide, n):
    add_text(slide, str(n), W - Inches(0.6), H - Inches(0.4), Inches(0.5), Inches(0.3),
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def logo(slide, dark=False):
    c = WHITE if dark else BRAND_DEEP
    add_text(slide, "AltHealth", Inches(0.4), Inches(0.22), Inches(2), Inches(0.45),
             size=14, bold=True, color=c, font="Fraunces")

def table_slide(slide, headers, rows, x, y, w, col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h  = Inches(0.42)
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, row_h * n_rows).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    def cell_text(cell, text, bold=False, color=INK, size=12, align=PP_ALIGN.LEFT, bg=None):
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.background()
        tf = cell.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name  = "Manrope"
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color

    for ci, h in enumerate(headers):
        cell_text(tbl.cell(0, ci), h, bold=True, color=WHITE, size=11,
                  align=PP_ALIGN.CENTER, bg=BRAND)

    for ri, row in enumerate(rows):
        bg = SOFT_BLUE if ri % 2 == 0 else CARD_BG
        for ci, val in enumerate(row):
            cell_text(tbl.cell(ri + 1, ci), val, color=INK, size=11,
                      align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT, bg=bg)

# ══════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════

# ── 1. TITLE ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
# gold left bar
add_rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
# tagline chip
chip = add_rect(s, Inches(0.6), Inches(1.1), Inches(3.5), Inches(0.38), fill=RGBColor(0x1e,0x4a,0xd4))
add_text(s, "VENTURE PARTNER OPPORTUNITY", Inches(0.7), Inches(1.12), Inches(3.3), Inches(0.34),
         size=9, bold=True, color=WHITE, font="Manrope")
add_text(s, "AltHealth\nVenture Network", Inches(0.6), Inches(1.7), Inches(8), Inches(2.2),
         size=52, bold=True, color=WHITE, font="Fraunces")
add_text(s, "Own the wellness commerce infrastructure in your market.\nRevenue share · Equity · Venture-style upside.",
         Inches(0.6), Inches(3.9), Inches(7.5), Inches(1),
         size=17, color=RGBColor(0xb8,0xcc,0xff), font="Manrope")
add_text(s, "althealth.me/venture  ·  support@althealth.me",
         Inches(0.6), H - Inches(0.65), Inches(6), Inches(0.4),
         size=11, color=RGBColor(0x88,0xa0,0xd0), font="Manrope", italic=True)
accent_bar(s)

# ── 2. THE PROBLEM ────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 2)
eyebrow(s, "The Problem", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "The wellness economy is broken\nat the distribution layer.", Inches(0.6), Inches(1.25), Inches(7.5), size=34)
items = [
    "Brands spend 40–60% of revenue on ads with declining returns",
    "Practitioners give away billions in purchasing influence — for free",
    "Patients buy from Amazon based on reviews instead of their doctor's advice",
    "No infrastructure connects trust → recommendation → purchase at scale",
]
bullet_box(s, items, Inches(0.6), Inches(3.0), Inches(6.5), Inches(2.5), size=15)
# stat callout
add_rect(s, Inches(8.4), Inches(1.4), Inches(4.3), Inches(4.6), fill=SOFT_BLUE, line=BRAND)
add_text(s, "$10T", Inches(8.6), Inches(1.8), Inches(3.9), Inches(1.2),
         size=64, bold=True, color=BRAND, font="Fraunces", align=PP_ALIGN.CENTER)
add_text(s, "Global wellness economy\nby 2030", Inches(8.6), Inches(3.1), Inches(3.9), Inches(0.8),
         size=14, color=MUTED, font="Manrope", align=PP_ALIGN.CENTER)
add_text(s, "1M+\nlicensed US practitioners", Inches(8.6), Inches(4.1), Inches(3.9), Inches(0.75),
         size=13, color=BRAND_DEEP, font="Manrope", align=PP_ALIGN.CENTER, bold=True)
add_text(s, "5–10×\nbetter conversion vs. digital ads", Inches(8.6), Inches(5.0), Inches(3.9), Inches(0.75),
         size=13, color=BRAND_DEEP, font="Manrope", align=PP_ALIGN.CENTER, bold=True)

# ── 3. THE SOLUTION ───────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 3)
eyebrow(s, "The Solution", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "AltHealth is the commerce OS\nfor practitioner-led wellness.", Inches(0.6), Inches(1.25), Inches(12), size=34)
# three columns
for i, (icon, title, desc, col) in enumerate([
    ("🏥", "Brands", "Authentic distribution with full attribution. No wasted ad spend. Real ROI.", ACCENT),
    ("👩‍⚕️", "Practitioners", "Passive income from recommendations already being made. No inventory. No ads.", BRAND),
    ("🧑‍🤝‍🧑", "Patients", "Buy what their trusted provider actually recommends — curated, transparent, trusted.", BRAND_DEEP),
]):
    cx = Inches(0.5 + i * 4.25)
    cw = Inches(4.0)
    add_rect(s, cx, Inches(2.7), cw, Inches(3.8), fill=CARD_BG, line=RGBColor(0xcc,0xd5,0xf0))
    add_text(s, icon, cx + Inches(0.2), Inches(2.9), Inches(0.8), Inches(0.6), size=28)
    add_text(s, title, cx + Inches(0.2), Inches(3.5), cw - Inches(0.4), Inches(0.5),
             size=17, bold=True, color=col, font="Fraunces")
    add_text(s, desc, cx + Inches(0.2), Inches(4.1), cw - Inches(0.4), Inches(1.8),
             size=13, color=MUTED, font="Manrope")

# ── 4. HOW IT WORKS ───────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 4)
eyebrow(s, "How It Works", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "Trust → Recommendation → Purchase.\nFully tracked and attributed.", Inches(0.6), Inches(1.25), Inches(12), size=32)
steps = [
    ("01", "Brand lists products", "Sets commission rates (10–25%), builds collections, joins campaigns."),
    ("02", "Practitioner curates", "Builds storefront, shares QR code or link with patients at point of care."),
    ("03", "Patient purchases", "Buys from the practitioner's storefront — familiar, trusted, frictionless."),
    ("04", "Everyone earns", "Brand pays commission → platform fee → practitioner earns. All attributed."),
]
for i, (num, title, desc) in enumerate(steps):
    cx = Inches(0.5 + i * 3.15)
    cw = Inches(2.95)
    add_rect(s, cx, Inches(2.85), cw, Inches(3.5), fill=CARD_BG, line=RGBColor(0xcc,0xd5,0xf0))
    add_text(s, num, cx + Inches(0.2), Inches(3.0), Inches(1), Inches(0.7),
             size=28, bold=True, color=SOFT_BLUE, font="Fraunces")
    add_text(s, title, cx + Inches(0.2), Inches(3.7), cw - Inches(0.3), Inches(0.5),
             size=14, bold=True, color=BRAND_DEEP, font="Fraunces")
    add_text(s, desc, cx + Inches(0.2), Inches(4.25), cw - Inches(0.3), Inches(1.7),
             size=12, color=MUTED, font="Manrope")

# ── 5. REVENUE MODEL ──────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 5)
eyebrow(s, "Business Model", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "Three recurring revenue streams.\nOperators share in all of them.", Inches(0.6), Inches(1.25), Inches(12), size=32)
table_slide(s,
    ["Revenue Stream", "Source", "Price Point", "Operator Share"],
    [
        ["Brand Subscriptions", "Monthly SaaS fee per brand", "$99 – $999 / mo", "Per agreement"],
        ["Practitioner Pro", "Monthly membership per practitioner", "$99 / mo", "Per agreement"],
        ["Platform Transaction Fees", "% of GMV through marketplace", "% of each sale", "Per agreement"],
    ],
    Inches(0.6), Inches(3.1), Inches(12.1),
    col_widths=[Inches(3.2), Inches(3.4), Inches(2.5), Inches(3.0)]
)
add_text(s, "Venture Partners earn 50% revenue share across all streams within their market.",
         Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
         size=13, bold=True, color=BRAND, font="Manrope", italic=True)

# ── 6. VENTURE PARTNER OPPORTUNITY ───────────────────────────
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
logo(s, dark=True)
accent_bar(s)
add_text(s, "6", W - Inches(0.6), H - Inches(0.4), Inches(0.5), Inches(0.3),
         size=9, color=RGBColor(0x44,0x55,0x88), align=PP_ALIGN.RIGHT)
eyebrow(s, "The Opportunity", Inches(0.6), Inches(0.9), Inches(8))
add_text(s, "We're not looking for investors.\nWe're looking for co-owners.", Inches(0.6), Inches(1.25), Inches(8.5), Inches(1.8),
         size=38, bold=True, color=WHITE, font="Fraunces")
perks = [
    "50% revenue share on all activity in your market",
    "Profit share on market P&L",
    "Equity grant commensurate with your buy-in level",
    "Sub-operator rights — build your own operator network",
    "Strategic vertical or regional market rights",
    "Dedicated support team + AltHealth Commerce OS",
]
bullet_box(s, perks, Inches(0.6), Inches(3.2), Inches(7), Inches(3.5),
           size=15, color=RGBColor(0xb8,0xcc,0xff))
# right callout
add_rect(s, Inches(8.6), Inches(1.4), Inches(4.1), Inches(5.0),
         fill=RGBColor(0x1e,0x4a,0xd4), line=RGBColor(0x44,0x77,0xff))
add_text(s, "Your market.\nYour brand.\nYour upside.", Inches(8.8), Inches(1.8), Inches(3.7), Inches(2.0),
         size=22, bold=True, color=WHITE, font="Fraunces")
add_text(s, "AltHealth provides the infrastructure. You own the market.", Inches(8.8), Inches(3.9), Inches(3.7), Inches(1.2),
         size=12, color=RGBColor(0xb8,0xcc,0xff), font="Manrope")

# ── 7. INVESTMENT TIERS ───────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 7)
eyebrow(s, "Investment Tiers", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "Three levels of partnership.\nAll receive 50% revenue share + equity.", Inches(0.6), Inches(1.25), Inches(12), size=32)
tier_data = [
    ("$25K", "Entry", SOFT_BLUE, BRAND, [
        "Equity grant (entry level)",
        "Regional market rights",
        "50% revenue share",
        "Profit share on market P&L",
        "Dedicated support team",
        "Sub-operator rights",
    ]),
    ("$50K", "Growth", RGBColor(0xd0,0xdd,0xff), BRAND, [
        "Enhanced equity grant",
        "Multi-market rights",
        "50% revenue share",
        "Profit share on market P&L",
        "Dedicated support team",
        "Priority deal flow",
    ]),
    ("$100K", "Strategic", RGBColor(0x1a,0x3d,0xb8), WHITE, [
        "Maximum equity grant",
        "National / vertical rights",
        "50% revenue share",
        "Profit share on market P&L",
        "Board observer seat",
        "JV & expansion rights",
    ]),
]
for i, (amount, label, bg, text_col, features) in enumerate(tier_data):
    cx = Inches(0.5 + i * 4.25)
    cw = Inches(4.05)
    add_rect(s, cx, Inches(2.75), cw, Inches(4.2), fill=bg, line=BRAND)
    add_text(s, amount, cx + Inches(0.25), Inches(2.95), cw - Inches(0.4), Inches(1.0),
             size=38, bold=True, color=text_col, font="Fraunces")
    add_text(s, label + " Partner", cx + Inches(0.25), Inches(3.85), cw - Inches(0.4), Inches(0.4),
             size=12, bold=True, color=text_col if bg != SOFT_BLUE else BRAND_DEEP, font="Manrope")
    for j, feat in enumerate(features):
        add_text(s, f"✓  {feat}", cx + Inches(0.25), Inches(4.35 + j * 0.33), cw - Inches(0.3), Inches(0.32),
                 size=11, color=text_col if i == 2 else MUTED, font="Manrope")

# ── 8. OPERATOR TIERS OVERVIEW ────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 8)
eyebrow(s, "Full Operator Tier Structure", Inches(0.6), Inches(0.9), Inches(9))
heading(s, "Start anywhere. Scale to Venture Partner.", Inches(0.6), Inches(1.25), Inches(12), size=32)
table_slide(s,
    ["Tier", "Buy-In", "Rev Share", "Equity", "Market Rights", "Support"],
    [
        ["🌱 Business Connector", "Free", "5–10%", "—", "None", "Self-serve"],
        ["🏗️ Market Operator", "$10K setup + $1K/mo", "30%", "—", "Conditional", "Group coaching"],
        ["🚀 Venture Partner — Entry", "$25K", "50%", "✓ Entry grant", "Regional", "Dedicated team"],
        ["🚀 Venture Partner — Growth", "$50K", "50%", "✓ Enhanced grant", "Multi-market", "Dedicated team"],
        ["🚀 Venture Partner — Strategic", "$100K", "50%", "✓ Maximum grant", "National/Vertical", "Dedicated + Board"],
    ],
    Inches(0.5), Inches(2.85), Inches(12.3),
    col_widths=[Inches(2.6), Inches(2.0), Inches(1.3), Inches(1.8), Inches(1.9), Inches(2.7)]
)

# ── 9. OPEN MARKETS ───────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 9)
eyebrow(s, "Open Markets", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "16+ open verticals. First mover wins the market.", Inches(0.6), Inches(1.25), Inches(9), size=32)
markets = [
    "Functional Medicine", "Women's Hormones", "Gut Health", "Longevity",
    "Mental Wellness", "Sports Recovery", "Biohacking", "Fertility",
    "Sleep Health", "Skin & Dermatology", "Ayurveda", "Naturopaths",
    "Yoga Studios", "Corporate Wellness", "Pet Wellness", "Chiropractors",
]
cols = 4
for i, m in enumerate(markets):
    col = i % cols
    row = i // cols
    cx = Inches(0.5 + col * 3.2)
    cy = Inches(2.9 + row * 0.72)
    add_rect(s, cx, cy, Inches(3.0), Inches(0.55),
             fill=SOFT_BLUE, line=RGBColor(0xb0,0xc4,0xf0))
    add_text(s, m, cx + Inches(0.15), cy + Inches(0.08), Inches(2.7), Inches(0.38),
             size=12, color=BRAND_DEEP, font="Manrope", bold=True)
add_text(s, "Exclusivity tied to performance milestones. Many categories still available.",
         Inches(0.6), Inches(6.85), Inches(10), Inches(0.4),
         size=11, color=MUTED, font="Manrope", italic=True)

# ── 10. TRACTION ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 10)
eyebrow(s, "Traction", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "Live platform. Founding cohorts open.\nEarly-mover economics available now.", Inches(0.6), Inches(1.25), Inches(12), size=32)
stats = [
    ("Live", "Platform at althealth.me"),
    ("25", "Founding brand spots"),
    ("44", "Founding practitioner spots"),
    ("3", "Operator tiers active"),
]
for i, (val, label) in enumerate(stats):
    cx = Inches(0.5 + i * 3.2)
    add_rect(s, cx, Inches(3.0), Inches(3.0), Inches(2.6), fill=SOFT_BLUE, line=BRAND)
    add_text(s, val, cx + Inches(0.2), Inches(3.2), Inches(2.6), Inches(1.0),
             size=44, bold=True, color=BRAND, font="Fraunces", align=PP_ALIGN.CENTER)
    add_text(s, label, cx + Inches(0.2), Inches(4.3), Inches(2.6), Inches(0.6),
             size=13, color=MUTED, font="Manrope", align=PP_ALIGN.CENTER)
add_text(s, "Founding partners receive the most favorable equity terms — terms improve as the platform scales.",
         Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
         size=13, bold=True, color=BRAND, font="Manrope", italic=True)

# ── 11. WHY NOW ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
light_slide_bg(s)
logo(s)
accent_bar(s)
slide_number(s, 11)
eyebrow(s, "Why Now", Inches(0.6), Inches(0.9), Inches(6))
heading(s, "Five tailwinds converging.", Inches(0.6), Inches(1.25), Inches(12), size=34)
reasons = [
    ("📈", "Wellness spending accelerating", "Post-pandemic health awareness has permanently elevated wellness spending across all demographics."),
    ("💰", "Brands pulling back on Meta/Google", "Rising CPMs and declining ROAS are forcing brands to find alternative, trust-based distribution channels."),
    ("👩‍⚕️", "Practitioners want new income streams", "Declining reimbursements are pushing practitioners to diversify — AltHealth is the natural fit."),
    ("🤖", "AI making curation scalable", "Personalized product recommendations can now be automated at scale, improving conversion for everyone."),
    ("⏱️", "Pre-scale = best economics", "AltHealth is pre-institutional. Founding Venture Partners get equity terms that won't be available later."),
]
for i, (icon, title, desc) in enumerate(reasons):
    cy = Inches(2.5 + i * 0.85)
    add_text(s, icon, Inches(0.6), cy, Inches(0.6), Inches(0.5), size=18)
    add_text(s, title, Inches(1.25), cy + Inches(0.03), Inches(3.5), Inches(0.4),
             size=13, bold=True, color=INK, font="Manrope")
    add_text(s, desc, Inches(4.9), cy + Inches(0.03), Inches(7.8), Inches(0.4),
             size=12, color=MUTED, font="Manrope")
    if i < 4:
        add_rect(s, Inches(0.6), cy + Inches(0.75), Inches(12.1), Inches(0.01),
                 fill=RGBColor(0xcc,0xd5,0xe8))

# ── 12. THE ASK ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
dark_slide_bg(s)
add_rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
logo(s, dark=True)
accent_bar(s)
add_text(s, "12", W - Inches(0.6), H - Inches(0.4), Inches(0.5), Inches(0.3),
         size=9, color=RGBColor(0x44,0x55,0x88), align=PP_ALIGN.RIGHT)
eyebrow(s, "The Ask", Inches(0.6), Inches(0.9), Inches(8))
add_text(s, "Join the founding cohort\nof Venture Partners.", Inches(0.6), Inches(1.25), Inches(9), Inches(2.0),
         size=42, bold=True, color=WHITE, font="Fraunces")
add_text(s, "Spots are limited. Market rights are first-come by vertical and region.\nFounding partners receive the most favorable equity terms available.",
         Inches(0.6), Inches(3.4), Inches(8.5), Inches(0.9),
         size=15, color=RGBColor(0xb8,0xcc,0xff), font="Manrope")
for i, (level, amount) in enumerate([("Entry", "$25K"), ("Growth", "$50K"), ("Strategic", "$100K")]):
    cx = Inches(0.6 + i * 2.8)
    add_rect(s, cx, Inches(4.5), Inches(2.6), Inches(1.4),
             fill=RGBColor(0x1e,0x4a,0xd4), line=RGBColor(0x44,0x77,0xff))
    add_text(s, amount, cx + Inches(0.15), Inches(4.65), Inches(2.3), Inches(0.6),
             size=26, bold=True, color=WHITE, font="Fraunces", align=PP_ALIGN.CENTER)
    add_text(s, level + " Partner", cx + Inches(0.15), Inches(5.25), Inches(2.3), Inches(0.4),
             size=11, color=RGBColor(0xb8,0xcc,0xff), font="Manrope", align=PP_ALIGN.CENTER)
add_text(s, "Apply:  althealth.me/venture", Inches(0.6), Inches(6.2), Inches(7), Inches(0.45),
         size=15, bold=True, color=WHITE, font="Manrope")
add_text(s, "Contact:  support@althealth.me", Inches(0.6), Inches(6.65), Inches(7), Inches(0.4),
         size=13, color=RGBColor(0xb8,0xcc,0xff), font="Manrope")

# ── Save ──────────────────────────────────────────────────────
out = "/home/user/althealth-web/AltHealth_VenturePartner_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
